from flask import Flask,render_template,url_for,redirect,request,jsonify
from flask_mail import Mail,Message 
import Text_summarization
from werkzeug.utils import secure_filename
from allennlp.predictors.predictor import Predictor
from summarizer import Summarizer
import os
import re
import docx
import PyPDF2 
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import json

# Make a regular expression 
# for validating an Email 
regex = "^[a-z0-9]+[\\._]?[a-z0-9]+[@]\\w+[.]\\w{2,3}$"

#app configurations
application=app=Flask(__name__)
app.config['UPLOAD_FOLDER']= 'uploads/'
app.config['MAIL_SERVER']='smtp.gmail.com'
app.config['MAIL_PORT']=587
app.config['MAIL_USE_TLS']=True
app.config['MAIL_USE_SSL']=False
app.config['MAIL_USERNAME']='aadarshgupta875@gmail.com'
app.config['MAIL_PASSWORD']=os.environ.get('password')
app.config['MAIL_DEFAULT_SENDER']='aadarshgupta875@gmail.com'

mail=Mail(app)

text=''
text_summary=''
pred_summary=''
name=''

predictor = Predictor.from_path("https://storage.googleapis.com/allennlp-public-models/bidaf-model-2020.03.19.tar.gz")


@app.route('/',methods=["GET","POST"])
def dashboard():
    return render_template('index.html')



def readfile(flag):
	global text
	global text_summary
	global name
	File=''

	#loads file from 
	if(flag==0):
		File=request.files['filename']
	else:
		File=request.files['file_summary']


	print(File)

	#contains filename
	fn=os.path.basename(File.filename)
	fn=fn.replace(" ", "_")
	print(fn)

	#saving file in specified directory
	File.save(os.path.join(app.config['UPLOAD_FOLDER'],secure_filename(File.filename)))

	#contains path where file is saved
	File=str(os.path.join(app.config['UPLOAD_FOLDER'])+fn)

	#dividing the file location into its filename and extension
	name, ext =os.path.splitext(File)

	#checking whether the file is videofile or not 
	allowed_extensions=[".doc",".docx",".pdf",".rtf",".pptx",".txt"]
	if (ext.lower() in allowed_extensions)==False:
		return "extension_error"

	if(ext.lower()=='.txt'):
		try:
			f = open(File, "r")
			data = f.read()
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.docx' or ext.lower()=='.doc'):
		try:
			doc = docx.Document(File)  # Creating word reader object.
			data = ""
			fullText = []
			for para in doc.paragraphs:
				fullText.append(para.text)
				data = '\n'.join(fullText)
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.pptx'):
		try:
			prs = Presentation(File)
			data = ""
			fullText = []
			for slide in prs.slides:
				for shape in slide.shapes:
					if hasattr(shape, "text"):
						fullText.append(shape.text)
				data = '\n'.join(fullText)
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.pdf'):
		try:
			pdfFileObj = open(File, mode='rb')    # creating a pdf file object 
			pdfReader = PyPDF2.PdfFileReader(pdfFileObj)    # creating a pdf reader object
			number_of_pages=pdfReader.numPages    #counting number of pages
			data=""
			fullText=[]
			for i in range(0,number_of_pages): 
				pageObj = pdfReader.getPage(i)    # a page object
				fullText.append(pageObj.extractText())   # extracting text from pdf 
			data=' '.join(fullText)
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.rtf'):
		try:
			f = open(File, 'r')
			rtf_text=f.read()
			data = rtf_to_text(rtf_text)
		except Exception:
			return "Something_went_wrong"
			

	# flag=0 -> Original_File
	# flag=1 -> Summary_File
	if(flag==0):
		text=data
	else:
		text_summary=data
	return "successfully_read"


@app.route('/email',methods=["GET","POST"])
def email():
	#retrieves email ID from html
	global pred_summary
	global name

	recipient_id=str(request.form.get("email"))

	#if invalid email is provided 
	if not(re.search(regex,recipient_id)):
		return render_template("summarize.html",text="Invalid Email Id")

	name=name[8:]
	print(name)
	#Summary file name
	Summary_file="summary_"+name+".txt"
	file1 = open(Summary_file, "w+")  
	file1.write(pred_summary) 
	file1.close()

	try:
		#Message Content that is to be emailed 
		msg=Message(
			subject='Summary file',
			recipients=[recipient_id],
			body='Below is your required summary'
			)

		#attaching summary file
		with app.open_resource(Summary_file) as output_file:
			msg.attach(Summary_file,'text/plain',output_file.read())

		#sending mail
		mail.send(msg)

		return render_template("summarize.html",text="Successfully mailed summary to your provided email id")
		
	except Exception:
		return render_template("summarize.html",text="Sorry!!! Unable to send the email")


@app.route('/summarize',methods=["GET","POST"])
def summarize():
	global pred_summary

	if request.method=="POST":
		radio_butt=request.form.get("optradio")
		print(radio_butt)
		output=readfile(0)
		if(output=="extension_error"):
			return render_template("summarize.html",summary_output="Not a text file")

		if(output=="Something_went_wrong"):
			return render_template("summarize.html",summary_output="Sorry!!! Something_went_wrong")

		if(radio_butt=="bert"):
			model = Summarizer()
			pred_summary = model(text)
			return render_template("summarize.html",summary_output=pred_summary)
		else:
			pred_summary=Text_summarization.get_data(text,'',0)
			return render_template("summarize.html",summary_output=pred_summary)
			
	return render_template('summarize.html')



@app.route('/analyse',methods=["GET","POST"])
def analyse():
	global text
	global text_summary
	if request.method=="POST":

		output_text=readfile(0)
		output_summary=readfile(1)

		#if original file is not a text file
		if(output_text=="extension_error"):
			return render_template("analyse.html",analyse_output="Original file is not a text file")

		if(output_summary=="extension_error"):
			return render_template("analyse.html",analyse_output="summary file is not a text file")

		if(output_text=="Something_went_wrong" or output_summary=="Something_went_wrong"):
			return render_template("analyse.html",analyse_output="Sorry!!! Something_went_wrong")
		
		accuracy=Text_summarization.get_data(text,text_summary,1)
		return render_template("analyse.html",analyse_output=accuracy)

	return render_template('analyse.html')



@app.route('/quesans',methods=["GET","POST"])
def quesans():
	return render_template('quesans.html')



@app.route('/file',methods=["POST"])
def file():
	global text

	output=readfile(0)

	if(output=="extension_error"):
		return render_template("quesans.html",qa_output="Not a text file")

	if(output=="successfully_read"):
		return render_template("quesans.html",qa_output="File uploaded successfully")

	if(output=="Something_went_wrong"):
		return render_template("quesans.html",qa_output="Sorry!!! Something_went_wrong")



@app.route('/question',methods=["GET","POST"])
def question():
	global text
	question=str(request.form.get('comment'))
	try:
		result=predictor.predict(
		  	passage=text,
	        question=question
		)
		answer=result['best_span_str']
		return render_template("quesans.html",question=question,answer=answer)
	except Exception:
		return render_template("quesans.html",answer="Sorry!!! Currently application is out of service")


@app.route('/MailMe', methods=['POST', 'GET'])
def MailMe():
	data= request.get_json(force = True)
	name=data.get("name")
	email=data.get("email")
	phone=data.get("phone")
	message=data.get("message")
	try:
		msg=Message(
			subject='Text Summarization Website Contact Form: %s'%(name),
			recipients=['aadarshgupta875@gmail.com'],
			body="You have received a new message from your website contact form.\n\n. Here are the details:\n\nName: {}\n\nEmail: {}\n\nPhone: {}\n\nMessage:\n{}".format(name,email,phone,message)
		)
		mail.send(msg)
		return json.dumps({'success':True})
	except Exception:
		return json.dumps({'error':True})


if __name__=="__main__":
	app.run()
